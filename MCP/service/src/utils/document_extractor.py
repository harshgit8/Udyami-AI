"""
Document text extraction utilities for local mode
"""
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


class DocumentExtractor:
    """Extract text content from various document formats"""
    
    @staticmethod
    def extract_text(file_path: Path) -> Optional[str]:
        """
        Extract text from a document file
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text content or None if extraction fails
        """
        ext = file_path.suffix.lower()
        
        try:
            if ext == '.pdf':
                return DocumentExtractor._extract_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                return DocumentExtractor._extract_docx(file_path)
            elif ext in ['.pptx', '.ppt']:
                return DocumentExtractor._extract_pptx(file_path)
            elif ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html']:
                return DocumentExtractor._extract_text_file(file_path)
            else:
                logger.warning(f"Unsupported file type: {ext}")
                return None
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return None
    
    @staticmethod
    def _extract_pdf(file_path: Path) -> str:
        """Extract text from PDF file"""
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(str(file_path))
            text_parts = []
            
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text.strip():
                        text_parts.append(f"--- Page {page_num} ---\n{text}")
                except Exception as e:
                    logger.warning(f"Error extracting page {page_num}: {e}")
                    continue
            
            full_text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from PDF with {len(reader.pages)} pages")
            return full_text
        except ImportError:
            logger.error("pypdf not installed. Install with: pip install pypdf")
            return None
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            return None
    
    @staticmethod
    def _extract_docx(file_path: Path) -> str:
        """Extract text from DOCX file"""
        try:
            from docx import Document
            
            doc = Document(str(file_path))
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            full_text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from DOCX")
            return full_text
        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return None
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            return None
    
    @staticmethod
    def _extract_pptx(file_path: Path) -> str:
        """Extract text from PPTX file"""
        try:
            from pptx import Presentation
            
            prs = Presentation(str(file_path))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
            
            full_text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from PPTX with {len(prs.slides)} slides")
            return full_text
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return None
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            return None
    
    @staticmethod
    def _extract_text_file(file_path: Path) -> str:
        """Extract text from plain text files"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    logger.info(f"Extracted {len(text)} characters from text file using {encoding}")
                    return text
                except UnicodeDecodeError:
                    continue
            
            logger.error(f"Could not decode text file with any encoding")
            return None
        except Exception as e:
            logger.error(f"Text file extraction error: {e}")
            return None
    
    @staticmethod
    def create_summary(text: str, max_length: int = 500) -> str:
        """Create a simple summary from extracted text"""
        if not text:
            return "No content extracted"
        
        # Clean up text
        text = " ".join(text.split())
        
        if len(text) <= max_length:
            return text
        
        # Return first max_length characters with ellipsis
        return text[:max_length].rsplit(' ', 1)[0] + "..."
